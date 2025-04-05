#!/usr/bin/env python
# coding: utf-8

# In[1]:


# imports

import os
import requests
from dotenv import load_dotenv
from bs4 import BeautifulSoup
from IPython.display import Markdown, display
from openai import OpenAI

from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api._errors import TranscriptsDisabled, NoTranscriptFound
import re

from urllib.parse import urlparse, parse_qs

from pptx import Presentation
import json
import pyautogui
import time
import pygetwindow as gw
import random


# In[2]:


#Configs:
EnvKey = 'OPENAI_API_KEY'
baseURL = 'https://api.deepseek.com'
LLMmodel = 'deepseek-chat'
LanguageCode = 'en'
ModelWindowSize = 128000
PresentationName = "QunatumTunneling"
DestinationDirectory = "D:/WorkSpace_AfterFormat/Generated_presentations"
DestinationDirectory = DestinationDirectory + "/" + PresentationName + ".pptx"
PreferedSlideCount = "10-12"


# In[3]:


WORD_LIMIT = (ModelWindowSize * 3)/16
CHAR_LIMIT = ModelWindowSize


# In[4]:


# Load environment variables in a file called .env

load_dotenv(override=True)
api_key = os.getenv('OPENAI_API_KEY')

openai = OpenAI()
openai = OpenAI(base_url=baseURL, api_key=api_key)

#response = openai.chat.completions.create(model=LLMmodel, messages=[{"role":"user", "content":message}])


# In[5]:


# Methods for Videos Path

def extract_video_id(url):
    parsed_url = urlparse(url)

    # Case 1: Standard YouTube URL (https://www.youtube.com/watch?v=VIDEO_ID)
    if parsed_url.hostname in ['www.youtube.com', 'youtube.com']:
        query_params = parse_qs(parsed_url.query)
        return query_params.get('v', [None])[0]

    # Case 2: Shortened YouTube URL (https://youtu.be/VIDEO_ID)
    elif parsed_url.hostname in ['youtu.be']:
        return parsed_url.path.lstrip('/')

    return None  # If no video ID is found




def fetch_transcript(video_id):
    try:
        # Try to get English transcript first
        transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=[LanguageCode])
    except NoTranscriptFound:
        try:
            languagecode = ''
            ytt_api = YouTubeTranscriptApi()
            transcript_list = ytt_api.list_transcripts(video_id)
            for item in transcript_list:
                if item.language_code:
                    languagecode = item.language_code
                    break
            for trans in transcript_list:
                if LanguageCode in trans.translation_languages and trans.is_translatable:
                    try:
                        transcript = transcript_list.find_manually_created_transcript([languagecode])
                    except NoTranscriptFound:
                        transcript = transcript_list.find_generated_transcript([languagecode])

                    translated_transcript = transcript.translate(LanguageCode).fetch()

                    for entry in translated_transcript:
                        transcript_text += " " +entry.text
                    return transcript_text

            # If English isn't available, get the default transcript in any available language
            if languagecode != '':
                transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=[languagecode])
                transcript_text = "***#Language is "+ languagecode + " #*** and here is the transcript: "
            else:
                return None
        except (NoTranscriptFound, TranscriptsDisabled):
            print("No transcript found.")
            return None

    # Combine all transcript text into a single string
    transcript_text = " ".join([entry['text'] for entry in transcript])
    return transcript_text


# In[6]:


# Website Path Classes

# Some websites need you to use proper headers when fetching them:
headers = {
 "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/117.0.0.0 Safari/537.36"
}

class Website:

    def __init__(self, url):
        """
        Create this Website object from the given url using the BeautifulSoup library
        """
        self.url = url
        response = requests.get(url, headers=headers)
        soup = BeautifulSoup(response.content, 'html.parser')
        self.title = soup.title.string if soup.title else "No title found"
        for irrelevant in soup.body(["script", "style", "img", "input"]):
            irrelevant.decompose()
        self.text = soup.body.get_text(separator="\n", strip=True)


# In[7]:


# Prompts to slides of title and bullet points for each discussed idea in the transcript 

def user_prompt_for_final(transcript_text):
    return transcript_text

def user_prompt_for_outline(transcript_text):
    user_prompt = f"You are looking at content from multiple sources"
    user_prompt += "\nPlease summarize and provide the outline for this, return this in outline format.  \n\n"
    user_prompt += transcript_text
    return user_prompt

system_prompt_outline = f"""
You are an AI assistant that summarizes and organizes long text into a high-level outline.  

### **Instructions:**
1. Analyze the given text and extract key topics.
2. Summarize each topic in a few bullet points.
3. Ensure that no duplicate ideas are included.

### **Output Format:**
{{
    "outline": [
        {{
            "topic": "Main Topic",
            "summary": [
                "Key Point 1",
                "Key Point 2"
            ]
        }}
    ]
}}
"""



system_prompt_final = """
You are an AI assistant that processes text from multiple sources (video transcripts, scraped websites).  
Each source discusses the same topic but may provide unique insights or perspectives.  

### **Your Tasks:**
1. Translate all content into """ + LanguageCode+""" if necessary.
2. Ignore text that might be navigation related
3. **Identify and remove redundant or repeated ideas.**
4. Compare information across sources and **retain only unique contributions, new facts, or alternative viewpoints.**
5. Organize the refined content into a structured PowerPoint outline.

You may receive an optional user instruction input that provides:
1. A preferred number of slides (as a single number or a range).
2. A suggested structure or flow for the presentation (e.g., agenda → introduction → key points → conclusion).

When such instructions are provided:
- Follow the suggested flow strictly unless it's clearly infeasible based on the source material.
- Ensure that the total number of slides is as close as possible to the requested number or within the provided range.
- Use the structure to group the main ideas logically. You may combine or reorganize the source content to fit this flow.

If no instruction is given, you may choose the most logical structure based on the content.

### **Output Format:**
Respond with a JSON object structured as follows:
{
    "slides": [
        {
            "title": "Slide Title",
            "points": [
                {
                    "text": "Main Point",
                    "subpoints": ["Subpoint 1", "Subpoint 2"]
                }
            ]
        }
    ]
}

- **Do NOT include repetitive points from different sources.**
- **Only keep unique insights or perspectives.**
- **Summarize each unique idea clearly before structuring slides.**
"""


def messages_for(transcript_text, state):
    if state == 'outline':
        return [
            {"role": "system", "content": system_prompt_outline},
            {"role": "user", "content": user_prompt_for_outline(transcript_text)}
        ]
    if state == 'final':
        return [
            {"role": "system", "content": system_prompt_final},
            {"role": "user", "content": user_prompt_for_final(transcript_text)}
        ]

def get_LLM_response(transcript_text, state):
    response = openai.chat.completions.create(
            model = LLMmodel,
            messages = messages_for(transcript_text, state)
    )
    return response.choices[0].message.content


# In[8]:


def GetTextFromURL(url):
    video_id = extract_video_id(url)
    if video_id is None:
        transcript_text = Website(url).text
    else:
        transcript_text = fetch_transcript(video_id)
    if(transcript_text):
        return transcript_text
    else:
        return None

def GetTranscriptsFromUrls(urlList):
    Transcripts = []
    for url in urlList:
        Transcripts.append(GetTextFromURL(url))
    return Transcripts


# In[9]:


def count_words(text):
    return len(text.split())

def chunk_text(transcripts):
    """Split text into chunks based on word and character limits."""
    chunks = []
    current_chunk = ""
    current_word_count = 0
    current_char_count = 0

    for transcript in transcripts:
        transcript_word_count = count_words(transcript)
        transcript_char_count = len(transcript)

        # Check if adding this transcript would exceed limits
        if (transcript_word_count > WORD_LIMIT or transcript_char_count > CHAR_LIMIT):
            chunks.append(transcript)
            current_chunk = ""
            current_word_count = 0
            current_char_count = 0
        elif (current_word_count + transcript_word_count > WORD_LIMIT or 
            current_char_count + transcript_char_count > CHAR_LIMIT):
            # Start a new chunk
            chunks.append(current_chunk.strip())
            current_chunk = ""
            current_word_count = 0
            current_char_count = 0
        else:
            # Append to current chunk
            current_chunk += " " + transcript
            current_word_count += transcript_word_count
            current_char_count += transcript_char_count

    # Append the last chunk
    if current_chunk:
        chunks.append(current_chunk.strip())

    return chunks

"""Recursively reduce multiple chunks into a single outline."""
def iterative_outline(chunks):
    res = chunks
    if len(res) == 1:
        state = 'final'
    else:
        state = 'outline'
    while state == 'outline':
        NewTranscripts = []
        for chunk in res:
            NewTranscripts.append(get_LLM_response(chunk, state))
        res = chunk_text(NewTranscripts)
        if len(res) == 1:
            state = 'final'

    return res[0]




# In[10]:


def GetProperJsonstring(jsonslides):
    start = jsonslides.find('{')
    end = jsonslides.rfind('}')
    if start != -1 and end != -1:
        properJsonString = jsonslides[start:end+1]
    return properJsonString


# In[11]:


# Function to add slides with points and subpoints
def add_slide(prs, title, points):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    slide.shapes.title.text = title

    # Set bullet points
    content = slide.placeholders[1].text_frame
    for point in points:
        p = content.add_paragraph()
        p.text = point["text"]

        # Add subpoints (indented)
        for sub in point.get("subpoints", []):
            sub_p = content.add_paragraph()
            sub_p.text = sub
            sub_p.level = 1  # Indent subpoints




# In[12]:


PreferedSlideCount = "10-12"
UserSourceNotes = """
Quantum tunneling is a quantum mechanical phenomenon where particles pass through a potential energy barrier that they classically shouldn’t be able to cross. According to classical physics, a particle with energy less than the barrier height would be reflected. However, due to the wave-like nature of particles in quantum mechanics, there is a non-zero probability that the particle's wavefunction extends beyond the barrier, allowing it to “tunnel” through.

This effect is crucial in many physical processes, such as nuclear fusion in stars (e.g., hydrogen nuclei overcoming repulsive forces), electron tunneling in semiconductors and tunnel diodes, and even the operation of scanning tunneling microscopes. The likelihood of tunneling depends on the barrier’s width and height—the thinner and lower the barrier, the higher the probability.

Quantum tunneling demonstrates the probabilistic nature of quantum mechanics and has practical applications in modern electronics, nuclear physics, and quantum computing.
"""
UserInstructions = """
INSTRUCTIONS:
I want the presentation to follow this flow:
1. Agenda
2. Introduction to the idea and concept of quantum tunneling
3. several slides including all ideas and details you got from sources and notes
4. Applications in Electronics and other fields
5. Benefits of quantum tunneling
6. Challenges of using this effect
7. Future outlook
"""
if PreferedSlideCount:
    UserInstructions += f"\nPlease aim for {PreferedSlideCount} slides."

if UserSourceNotes:
    UserInstructions = UserInstructions + "\n\nUSER NOTES:\n" + "Below is a summary or outline from the user. Treat it as an additional content source " + "to be merged with transcripts and scraped data:\n\n" + f"{UserSourceNotes}\n" + "\n---\n" + "The following content comes from multiple sources. Please deduplicate, integrate, and generate slides accordingly.\n"
else:
    UserInstructions = UserInstructions + "\n\nSOURCE INTEGRATION:\n" + "You are looking at content from multiple sources. Please deduplicate, extract key ideas, and structure it into slides following the flow above.\n"

print(UserInstructions)


# In[13]:


urlist = [
    'https://phys.libretexts.org/Bookshelves/University_Physics/University_Physics_(OpenStax)/University_Physics_III_-_Optics_and_Modern_Physics_(OpenStax)/07%3A_Quantum_Mechanics/7.07%3A_Quantum_Tunneling_of_Particles_through_Potential_Barriers',
         'https://en.wikipedia.org/wiki/Quantum_tunnelling', 'https://en.wikipedia.org/wiki/Tunnel_junction', 'https://physics.stackexchange.com/questions/466615/quantum-tunneling-in-zener-diodes', 
'https://www.youtube.com/watch?v=Vi0Jxg3t08E']
Transcripts = GetTranscriptsFromUrls(urlist)
finalScript = iterative_outline(Transcripts)
state = 'final'
#print(UserInstructions)
#print(type(finalScript))
#print(type(UserInstructions))
#print(finalScript)
script2 = UserInstructions + finalScript
#print(script2)
jsonslides = get_LLM_response(script2, state)
#print(jsonslides)
# Parse 
jstring = GetProperJsonstring(jsonslides)
slides_content = json.loads(jstring)

slidecount = len(slides_content["slides"])

# Create a PowerPoint presentation
prs = Presentation()


# Generate slides
for slide_data in slides_content["slides"]:
    add_slide(prs, slide_data["title"], slide_data["points"])

# Save the presentation
prs.save(DestinationDirectory)
print("Presentation created successfully!")



# In[17]:


def apply_design():
    """ Applies a random design to the current slide without reopening Designer """

    # Step 1: Move to the first design (10 'Up' presses)
    for _ in range(5):
        pyautogui.press('up')
        time.sleep(0.1)

    # Step 2: Randomly select a design (0-9 Down presses)
    num_down_presses = random.randint(0, 4)
    for _ in range(num_down_presses):
        pyautogui.press('down')
        time.sleep(0.2)

    # Step 3: Apply the selected design
    pyautogui.press('enter')
    time.sleep(1)
    pyautogui.press('esc')
    time.sleep(0.5)

def process_all_slides(num_slides):
    """ Iterates through all slides and applies a random design while keeping Designer open """

    time.sleep(10)
    # Open Designer once at the start
    pyautogui.hotkey('alt', 'g', 'd')
    time.sleep(5)

    for slide in range(num_slides):
        apply_design()

        # Step 4: Move to the next slide and wait for designer
        if slide < num_slides - 1:
            pyautogui.press('down')
            time.sleep(4)
            # Step 5: go to designer
            pyautogui.press('f6')
            time.sleep(0.5)

    pyautogui.hotkey('ctrl', 's')
    time.sleep(1)
    pyautogui.hotkey('alt', 'f4')

os.startfile(DestinationDirectory)
time.sleep(6)
WindowName = PresentationName + ".pptx" + " - PowerPoint"
gw.getWindowsWithTitle(WindowName)[0].activate()

# Adjust the number of slides as needed
process_all_slides(slidecount)


